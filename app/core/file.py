import requests
import os
import mimetypes
import fitz
import json

from docx import Document
from pptx import Presentation
from PIL import Image
import pytesseract


# =====================================
#  with open을 사용하여 파일을 읽을 때
# =====================================
"""
원래 파일을 읽거나 쓰면 무조건 마지막에 close() 호출을 해야 리소스 낭비가 없음.
python에서 제공하는 open을 사용하면 자동적으로 마지막에 close를 해줌.

파일을 읽을때는 '파일포인터'로 현재 위치를 관리함. 그래서 with문 중간이나 마지막에서 다시 처음부터 읽고 싶다면 seek(0)를 호출해서 포인터를 처음 위치로 이동한 후 다시 읽어야함.
> 파일이 열릴 때:
  - 기본적으로 파일 포인터는 파일의 시작(0번 바이트)에 위치함.
> 파일 읽기/쓰기 작업을 하면:
  - 파일 포인터가 자동으로 이동함.
  - 파일을 읽으면 읽은 바이트 수만큼 포인터가 이동함.

"""


def open_txt_file(file_path):
    with open(file_path, "r", encoding="utf-8") as r:
        data = r.read()
        print(data)


def open_docx_file(file_path):
    """
    document.part.rels 쓰면 이미지 있는 부분 이미지 추출 가능.

    for rel in doc.part.rels:
        if "image" in doc.part.rels[rel].target_ref:  # 이미지인지 확인
            image_data = doc.part.rels[rel].target_part.blob
            with open(f"extracted_image_{image_count}.png", "wb") as img_file:
                img_file.write(image_data)
            image_count += 1

    print(f"{image_count}개의 이미지를 추출했습니다.")

    """
    # file_path = "app/assets/김건희 도이치모터스.docx"
    doc = Document(file_path)

    for chunk in doc.paragraphs:
        # time.sleep(0.5)
        print(chunk.text)


def open_docx_file_with_image(file_path):
    """docx 파일 읽으면서 이미지 있는 부분은 테서렉트로 OCR"""
    doc = Document(file_path)

    output_folder = "app/assets/extracted_images"
    os.makedirs(output_folder, exist_ok=True)

    image_count = 0
    processed_texts = set()

    for block in doc.element.body:
        # 텍스트 추출 (다양한 태그에 포함된 텍스트를 처리)
        if block.tag.endswith("p"):  # 문단 (paragraph)
            # print(block.tag)
            for para in block.iter():
                # print(para.tag)
                if para.tag.endswith("t") and para.text:  # w:t (text) 태그
                    para_text = para.text.strip()
                    if para_text not in processed_texts:
                        processed_texts.add(para_text)
                        print(para_text)

                # 이미지가 있을 때, 해당 참조를 찾아서 이미지 추출
                elif para.tag.endswith("pic"):
                    for rel in doc.part.rels:
                        if "image" in doc.part.rels[rel].target_ref:
                            image_data = doc.part.rels[rel].target_part.blob
                            image_path = os.path.join(
                                output_folder, f"image_{image_count}.png"
                            )

                            # 이미지 저장
                            with open(image_path, "wb") as img_file:
                                img_file.write(image_data)

                            print(f"[이미지 저장됨: {image_path}]")

                            # OCR 실행 (이미지에서 텍스트 추출)
                            img = Image.open(image_path)
                            extracted_text = pytesseract.image_to_string(img)
                            print(f"[OCR 결과 - {image_path}]:\n{extracted_text}\n")

                            image_count += 1

    if image_count == 0:
        print("이미지가 없습니다.")
    else:
        print(f"{image_count}개의 이미지가 처리되었습니다.")


def open_pptx_file(file_path):
    # file_path = "app/assets/랭체인과 프롬프트.pptx"
    ppt = Presentation(file_path)

    for i, slide in enumerate(ppt.slides):
        print(f"슬라이드 [{i+1}]")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                print(shape.text)


def open_pdf_file(file_path):
    # file_path = "app/assets/랭체인 완벽입문 정리.pdf"
    pdf = fitz.open(file_path)

    for page_num in range(len(pdf)):
        page = pdf[page_num]
        text = page.get_text("text")

        print(f"📌 [페이지 {page_num + 1}]")
        print(text)
        print("-" * 50)

    pdf.close()


def open_json_file(file_path):
    # file_path = "app/assets/law_103962_20100331.json"
    with open(file_path, "r", encoding="utf-8") as f:
        json_data = json.load(f)
        print(json_data)


async def file_download(file_path: str, file_name: str):
    save_path = f"app/assets/D_{file_name}"
    os.makedirs(os.path.dirname(save_path), exist_ok=True)

    if file_path.startswith("http"):
        try:
            response = requests.get(url=file_path, stream=True)
            response.raise_for_status()
        except:
            raise ValueError(f"❌ 다운로드 실패: {response.status_code}")

        with open(save_path, "wb") as f:  # wb는 바이너리 쓰기 모드
            for chunk in response.iter_content(chunk_size=8192):  # 8KB 단위로 다운로드
                f.write(chunk)
            print(f"✅ PDF 파일 다운로드 완료: {save_path}")

    else:
        file_type = mimetypes.guess_type(file_name)
        file_extension = mimetypes.guess_extension(file_type[0])

        match file_extension:
            case ".pdf":
                open_pdf_file(file_path)
            case ".docx":
                open_docx_file(file_path)
            case ".pptx":
                open_pptx_file(file_path)
            case ".json":
                open_json_file(file_path)
            case _:
                raise ValueError(f"{file_name} 파일은 지원하는 확장자가 아닙니다.")


if __name__ == "__main__":
    # file_paths = [
    #     "https://www.w3.org/WAI/ER/tests/xhtml/testfiles/resources/pdf/dummy.pdf",
    #     "app/assets/랭체인 완벽입문 정리.pdf",
    #     "https://www.w3.org/TR/PNG/iso_8859-1.txt",
    #     "app/assets/김건희 도이치모터스.docx",
    #     "app/assets/참사 지원금 수령.docx",
    #     "app/assets/law_103962_20100331.json",
    #     "app/assets/랭체인과 프롬프트.pptx"
    # ]

    # for i, file_path in enumerate(file_paths):
    #     print(f"[{i+1}] {file_path}")

    file_path = "app/assets/조선비즈 화제성 기사 생성 이슈사항.docx"
    open_docx_file_with_image(file_path)
