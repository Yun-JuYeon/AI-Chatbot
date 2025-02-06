from ast import match_case
import requests
import os
import time
import mimetypes
import fitz
import json

from docx import Document
from pptx import Presentation


# =====================================
#  with open을 사용하여 파일을 읽을 때
# =====================================
"""
원래 파일을 읽거나 쓰면 무조건 마지막에 close() 호출을 해야 리소스 낭비가 없음.
python에서 제공하는 open을 사용하면 자동적으로 마지막에 close를 해줌.

파일을 읽을때는 '파일포인터'로 현재 위치를 관리함. 그래서 with문 중간이나 마지막에서 다시 처음부터 읽고 싶다면 seek(0)를 호출해서 포인터를 처음 위치로 이동한 후 다시 읽어야함.
> 파일이 열릴 때:
  - 기본적으로 파일 포인터는 파일의 시작(0번 바이트)에 위치함.
> 파일 읽기/쓰기 작업을 하면:
  - 파일 포인터가 자동으로 이동함.
  - 파일을 읽으면 읽은 바이트 수만큼 포인터가 이동함.

"""
def open_txt_file(file_path): 
    with open(file_path, "r", encoding="utf-8") as r:
        data = r.read()
        print(data)

def open_docx_file(file_path):
    # file_path = "assets/김건희 도이치모터스.docx"
    doc = Document(file_path)

    for chunk in doc.paragraphs:
        # time.sleep(0.5)
        print(chunk.text)

def open_pptx_file(file_path):
    # file_path = "assets/랭체인과 프롬프트.pptx"
    ppt = Presentation(file_path)

    for i, slide in enumerate(ppt.slides):
        print(f"슬라이드 [{i+1}]")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                print(shape.text)

def open_pdf_file(file_path):
    # file_path = "assets/랭체인 완벽입문 정리.pdf"
    pdf = fitz.open(file_path)

    for page_num in range(len(pdf)):
        page = pdf[page_num]
        text = page.get_text("text")

        print(f"📌 [페이지 {page_num + 1}]")
        print(text)
        print("-" * 50)
    
    pdf.close()

def open_json_file(file_path):
    # file_path = "assets/law_103962_20100331.json"
    with open(file_path, "r", encoding="utf-8") as f:
        json_data = json.load(f)
        print(json_data)
        

async def file_download(file_path:str, file_name:str):
    save_path = f"assets/D_{file_name}"
    os.makedirs(os.path.dirname(save_path), exist_ok=True)

    if file_path.startswith("http"):
        try:
            response = requests.get(url=file_path, stream=True)
            response.raise_for_status()
        except:
            raise ValueError (f"❌ 다운로드 실패: {response.status_code}")

        with open(save_path, "wb") as f:  # wb는 바이너리 쓰기 모드
            for chunk in response.iter_content(chunk_size=8192): # 8KB 단위로 다운로드
                f.write(chunk)
            print(f"✅ PDF 파일 다운로드 완료: {save_path}")
        
    else:
        file_type = mimetypes.guess_type(file_name)
        file_extension = mimetypes.guess_extension(file_type[0])

        match file_extension:
            case ".pdf":
                open_pdf_file(file_path)
            case ".docx":
                open_docx_file(file_path)
            case ".pptx":
                open_pptx_file(file_path)
            case ".json":
                open_json_file(file_path)
            case _ :
                raise ValueError (f"{file_name} 파일은 지원하는 확장자가 아닙니다.")


if __name__=="__main__":
    file_paths = [
        "https://www.w3.org/WAI/ER/tests/xhtml/testfiles/resources/pdf/dummy.pdf",
        "assets/랭체인 완벽입문 정리.pdf",
        "https://www.w3.org/TR/PNG/iso_8859-1.txt",
        "assets/김건희 도이치모터스.docx",
        "assets/참사 지원금 수령.docx",
        "assets/law_103962_20100331.json",
        "assets/랭체인과 프롬프트.pptx"
    ]
    
    for i, file_path in enumerate(file_paths):
        print(f"[{i+1}] {file_path}")
